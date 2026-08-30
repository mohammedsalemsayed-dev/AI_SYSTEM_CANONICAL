"""Renderers (MILESTONE_M_PLAN.md §2, §16).

Markdown + HTML are stdlib. DOCX + PPTX render via python-docx / python-pptx
(integrations, not builds — §16; `RendererUnavailable` if the package is missing).
PDF renders via the bundled dependency-free writer (`authoring/pdf_writer.py`,
Helvetica, no install). Binary renderers put bytes in `RenderedDoc.data`; text
renderers use `.text`.
"""

from __future__ import annotations

import html
import io
from dataclasses import dataclass, field
from typing import Protocol, runtime_checkable

from app.services.authoring.model import DocumentModel, Section, SlideDeck


class RendererUnavailable(RuntimeError):
    pass


@dataclass
class RenderedDoc:
    mime: str
    text: str = ""            # for text formats; a short preview for binary ones
    ext: str = "txt"
    data: bytes = field(default=b"", repr=False)  # populated for binary formats

    @property
    def is_binary(self) -> bool:
        return bool(self.data)

    def payload(self) -> bytes:
        return self.data if self.data else self.text.encode("utf-8")


@runtime_checkable
class Renderer(Protocol):
    def render(self, doc: DocumentModel) -> RenderedDoc: ...


class MarkdownRenderer:
    def render(self, doc: DocumentModel) -> RenderedDoc:
        nums = doc.numbered_citations()
        out: list[str] = []
        if doc.title:
            out.append(f"# {doc.title}\n")
        for sec in doc.sections:
            self._section(sec, out, nums, base=2)
        cits = doc.all_citations()
        if cits:
            out.append("\n## References\n")
            for i, c in enumerate(cits, 1):
                label = c.label or (f"{c.source}" if c.source else c.id)
                out.append(f"{i}. {label}" + (f" ({c.trust})" if c.trust not in ("workspace", "user") else ""))
        return RenderedDoc("text/markdown", "\n".join(out).rstrip() + "\n", "md")

    def _section(self, sec: Section, out: list[str], nums: dict[str, int], *, base: int) -> None:
        out.append(f"\n{'#' * min(base, 6)} {sec.title}\n")
        for b in sec.blocks:
            self._block(b, out, nums)
        for child in sec.children:
            self._section(child, out, nums, base=base + 1)

    def _block(self, b, out: list[str], nums: dict[str, int]) -> None:
        cite = self._cite(b.citation_ids, nums)
        if b.kind == "heading":
            out.append(f"{'#' * min(max(b.level, 1), 6)} {b.text}")
        elif b.kind == "paragraph":
            out.append((b.text + cite).strip())
        elif b.kind == "quote":
            out.append(f"> {b.text}{cite}")
        elif b.kind == "list":
            out.extend(f"- {it}" for it in b.items)
            if cite:
                out.append(f"  {cite.strip()}")
        elif b.kind == "code":
            out.append(f"```{b.lang}\n{b.text}\n```")
        elif b.kind == "table" and b.rows:
            head, *body = b.rows
            out.append("| " + " | ".join(head) + " |")
            out.append("| " + " | ".join("---" for _ in head) + " |")
            out.extend("| " + " | ".join(r) + " |" for r in body)
        elif b.kind == "figure":
            out.append(f"_Figure: {b.caption or b.text}_{cite}")
        out.append("")

    @staticmethod
    def _cite(ids: list[str], nums: dict[str, int]) -> str:
        ns = sorted({nums[i] for i in ids if i in nums})
        return f" [{', '.join(str(n) for n in ns)}]" if ns else ""


class HtmlRenderer:
    def render(self, doc: DocumentModel) -> RenderedDoc:
        nums = doc.numbered_citations()
        e = html.escape
        out = ["<article>"]
        if doc.title:
            out.append(f"<h1>{e(doc.title)}</h1>")
        for sec in doc.sections:
            self._section(sec, out, nums, e, level=2)
        cits = doc.all_citations()
        if cits:
            out.append('<section class="references"><h2>References</h2><ol>')
            for c in cits:
                out.append(f"<li>{e(c.label or c.source or c.id)}</li>")
            out.append("</ol></section>")
        out.append("</article>")
        return RenderedDoc("text/html", "\n".join(out), "html")

    def _section(self, sec, out, nums, e, *, level: int) -> None:
        out.append("<section>")
        out.append(f"<h{min(level, 6)}>{e(sec.title)}</h{min(level, 6)}>")
        for b in sec.blocks:
            self._block(b, out, nums, e)
        for child in sec.children:
            self._section(child, out, nums, e, level=level + 1)
        out.append("</section>")

    def _block(self, b, out, nums, e) -> None:
        cite = self._cite(b.citation_ids, nums)
        if b.kind == "paragraph":
            out.append(f"<p>{e(b.text)}{cite}</p>")
        elif b.kind == "quote":
            out.append(f"<blockquote>{e(b.text)}{cite}</blockquote>")
        elif b.kind == "list":
            out.append("<ul>" + "".join(f"<li>{e(it)}</li>" for it in b.items) + "</ul>")
        elif b.kind == "code":
            out.append(f"<pre><code>{e(b.text)}</code></pre>")
        elif b.kind == "table" and b.rows:
            head, *body = b.rows
            out.append("<table><thead><tr>" + "".join(f"<th>{e(c)}</th>" for c in head) + "</tr></thead><tbody>")
            out.extend("<tr>" + "".join(f"<td>{e(c)}</td>" for c in r) + "</tr>" for r in body)
            out.append("</tbody></table>")
        elif b.kind == "figure":
            out.append(f'<figure><figcaption>{e(b.caption or b.text)}</figcaption></figure>')

    @staticmethod
    def _cite(ids, nums) -> str:
        ns = sorted({nums[i] for i in ids if i in nums})
        return f' <sup>[{", ".join(str(n) for n in ns)}]</sup>' if ns else ""


def _plain_preview(doc: DocumentModel) -> str:
    """A short readable summary of a rendered binary doc, for the transcript."""
    lines = [f"# {doc.title}" if doc.title else "(untitled)"]
    for sec in doc.walk():
        lines.append(f"\n{'#' * min(sec.level + 1, 6)} {sec.title}")
        body = sec.text_body().strip()
        if body:
            lines.append(body[:400] + ("…" if len(body) > 400 else ""))
    cits = doc.all_citations()
    if cits:
        lines.append(f"\n{len(cits)} reference(s).")
    return "\n".join(lines).strip() + "\n"


def _num(ids, nums) -> str:
    ns = sorted({nums[i] for i in ids if i in nums})
    return f" [{', '.join(map(str, ns))}]" if ns else ""


class DocxRenderer:
    """Themed Word document via python-docx: an accent cover page, recoloured
    heading styles, an accent rule under H1s, a footer with page numbers."""

    def __init__(self, theme=None) -> None:
        from app.services.authoring.theme import get_theme
        self.theme = theme if theme is not None else get_theme(None)

    def render(self, doc: DocumentModel) -> RenderedDoc:
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
            from docx.shared import Pt, RGBColor
        except ImportError as e:  # pragma: no cover - dependency is declared
            raise RendererUnavailable("DOCX needs python-docx (`pip install python-docx`)") from e

        th = self.theme
        rgb = lambda c: RGBColor(*c)  # noqa: E731
        nums = doc.numbered_citations()
        d = Document()

        normal = d.styles["Normal"]
        normal.font.name = th.body_font
        normal.font.size = Pt(11)
        normal.font.color.rgb = rgb(th.ink)
        for lvl, size in ((1, 20), (2, 15), (3, 12.5)):
            try:
                st = d.styles[f"Heading {lvl}"]
                st.font.name = th.heading_font
                st.font.size = Pt(size)
                st.font.color.rgb = rgb(th.heading)
                st.font.bold = lvl == 1
            except KeyError:
                pass

        # --- cover page ---
        if doc.title:
            for _ in range(6):
                d.add_paragraph()
            p = d.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.LEFT
            r = p.add_run(doc.title)
            r.font.name = th.heading_font
            r.font.size = Pt(34)
            r.font.bold = True
            r.font.color.rgb = rgb(th.accent)
            rule = d.add_paragraph()
            rr = rule.add_run("—" * 18)
            rr.font.color.rgb = rgb(th.accent)
            meta = d.add_paragraph()
            mr = meta.add_run(f"{doc.kind.capitalize()}  ·  {len(list(doc.walk()))} sections")
            mr.font.size = Pt(11)
            mr.font.color.rgb = rgb(th.dim)
            d.add_page_break()

        def accent_rule(p) -> None:
            pPr = p._p.get_or_add_pPr()
            bdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "18")
            bottom.set(qn("w:color"), "%02X%02X%02X" % th.accent)
            bdr.append(bottom)
            pPr.append(bdr)

        def emit_section(sec: Section, base: int) -> None:
            h = d.add_heading(sec.title, level=min(base, 9))
            if base == 1:
                accent_rule(h)
            for b in sec.blocks:
                cite = _num(b.citation_ids, nums)
                if b.kind == "heading":
                    d.add_heading(b.text, level=min(max(b.level, 1), 9))
                elif b.kind == "paragraph":
                    d.add_paragraph((b.text + cite).strip())
                elif b.kind == "quote":
                    q = d.add_paragraph((b.text + cite).strip())
                    q.paragraph_format.left_indent = Pt(18)
                    for run in q.runs:
                        run.italic = True
                        run.font.color.rgb = rgb(th.dim)
                elif b.kind == "list":
                    for it in b.items:
                        d.add_paragraph(it, style="List Bullet")
                    if cite.strip():
                        d.add_paragraph(cite.strip())
                elif b.kind == "code":
                    p = d.add_paragraph()
                    run = p.add_run(b.text)
                    run.font.name = th.mono_font
                    run.font.size = Pt(9.5)
                elif b.kind == "table" and b.rows:
                    head, *body = b.rows
                    t = d.add_table(rows=1, cols=len(head))
                    t.style = "Light Grid Accent 1"
                    for i, c in enumerate(head):
                        t.rows[0].cells[i].text = str(c)
                    for r in body:
                        cells = t.add_row().cells
                        for i, c in enumerate(r[: len(head)]):
                            cells[i].text = str(c)
                elif b.kind == "figure":
                    d.add_paragraph(f"Figure: {b.caption or b.text}{cite}", style="Caption")
            for child in sec.children:
                emit_section(child, base + 1)

        for sec in doc.sections:
            emit_section(sec, 1)

        cits = doc.all_citations()
        if cits:
            h = d.add_heading("References", level=1)
            accent_rule(h)
            for i, c in enumerate(cits, 1):
                d.add_paragraph(f"{i}. {c.label or c.source or c.id}", style="List Number")

        # footer: title (left) + PAGE field (right)
        try:
            footer = d.sections[0].footer.paragraphs[0]
            footer.text = (doc.title or "")[:60] + "\t\t"
            for run in footer.runs:
                run.font.size = Pt(8)
                run.font.color.rgb = rgb(th.dim)
            fld = OxmlElement("w:fldSimple")
            fld.set(qn("w:instr"), "PAGE")
            r = OxmlElement("w:r")
            rpr = OxmlElement("w:rPr")
            sz = OxmlElement("w:sz"); sz.set(qn("w:val"), "16"); rpr.append(sz)
            r.append(rpr)
            fld.append(r)
            footer._p.append(fld)
        except Exception:  # noqa: BLE001 — footer is cosmetic
            pass

        buf = io.BytesIO()
        d.save(buf)
        return RenderedDoc(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            text=_plain_preview(doc), ext="docx", data=buf.getvalue(),
        )


import re as _re

_STAT_RE = _re.compile(
    r"^\s*[\$£€]?\d[\d.,]*\s*(%|x|×|k|m|bn|b|billion|million|percent|pt|pts|\+)?\b", _re.I
)

# --- auto visuals for slide decks (Gamma-style) ----------------------------- #
# When no images are attached, sprinkle a deterministic stock photo onto some
# standard slides. The seed is the slide topic, so re-rendering the same deck
# gives the same pictures. Offline / fetch failure -> a themed colour panel with
# a caption, so the deck is always valid. Opt out with NEXUS_DECK_IMAGES=0.
_PLACEHOLDER_EVERY = 3          # ~ every 3rd standard slide gets a visual
_PLACEHOLDER_TIMEOUT = 4.0      # seconds per fetch before falling back


def _img_slug(text: str) -> str:
    s = _re.sub(r"[^a-z0-9]+", "-", (text or "").lower()).strip("-")
    return (s or "abstract")[:48]


def _deck_images_enabled() -> bool:
    import os
    return os.environ.get("NEXUS_DECK_IMAGES", "1").strip().lower() not in ("0", "false", "no")


def _fetch_placeholder(topic: str, w: int = 1280, h: int = 860):
    """A deterministic stock photo for `topic` (picsum.photos), cached under the
    OS temp dir. Returns a local file path, or None when disabled / offline /
    the fetch fails — callers then draw a themed panel instead."""
    if not _deck_images_enabled():
        return None
    import hashlib
    import os
    import tempfile
    import urllib.request

    slug = _img_slug(topic)
    key = hashlib.sha1(f"{slug}-{w}x{h}".encode()).hexdigest()[:16]
    cache_dir = os.path.join(tempfile.gettempdir(), "nexus_deck_img")
    try:
        os.makedirs(cache_dir, exist_ok=True)
    except OSError:
        return None
    path = os.path.join(cache_dir, f"{key}.jpg")
    if os.path.isfile(path) and os.path.getsize(path) > 512:
        return path
    url = f"https://picsum.photos/seed/{slug}/{w}/{h}"
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "nexus-authoring"})
        with urllib.request.urlopen(req, timeout=_PLACEHOLDER_TIMEOUT) as r:  # noqa: S310
            data = r.read()
        if len(data) < 512:
            return None
        with open(path, "wb") as f:
            f.write(data)
        return path
    except Exception:  # noqa: BLE001 — offline is the expected non-fatal case
        return None


class PptxRenderer:
    """Themed deck via python-pptx with per-slide layout selection: a section with
    sub-sections becomes a full-bleed divider; a short stat-heavy slide becomes a
    big-number slide; a single long line becomes a centred quote; 5+ bullets go
    two-column; attached images get their own image+text layout. Otherwise a
    standard bullet slide. Coloured backgrounds, accent bar, footer + numbers."""

    def __init__(self, theme=None, images: list[str] | None = None) -> None:
        from app.services.authoring.theme import get_theme
        self.theme = theme if theme is not None else get_theme(None)
        self.images = [p for p in (images or []) if p]

    @staticmethod
    def _layout(sl, has_children: bool, has_image: bool) -> str:
        b = [x for x in (sl.bullets or []) if x.strip()]
        if has_children and not b:
            return "divider"
        if has_image:
            return "image"
        if len(b) == 1 and (len(b[0]) > 90 or b[0].lstrip().startswith(('"', "“"))):
            return "quote"
        stats = [x for x in b if _STAT_RE.match(x) and len(x) < 60]
        if b and len(b) <= 3 and len(stats) >= max(1, len(b) - 1):
            return "bignum"
        if len(b) >= 5:
            return "twocol"
        return "standard"

    def render(self, doc: DocumentModel) -> RenderedDoc:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as e:  # pragma: no cover
            raise RendererUnavailable("PPTX needs python-pptx (`pip install python-pptx`)") from e

        th = self.theme
        rgb = lambda c: RGBColor(*c)  # noqa: E731
        W, H = Inches(13.333), Inches(7.5)
        deck = SlideDeck.from_document(doc)
        prs = Presentation()
        prs.slide_width, prs.slide_height = W, H
        blank = prs.slide_layouts[6]

        def bg(slide) -> None:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb(th.bg)

        def textbox(slide, l, t, w, h, text, *, size, color, font, bold=False, align=PP_ALIGN.LEFT):
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            first = True
            for line in (text if isinstance(text, list) else [text]):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.text = line
                p.alignment = align
                for r in p.runs:
                    r.font.name = font
                    r.font.size = Pt(size)
                    r.font.bold = bold
                    r.font.color.rgb = rgb(color)
            return box

        def accent_line(slide, y) -> None:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), y, Inches(2.2), Pt(4))
            bar.fill.solid(); bar.fill.fore_color.rgb = rgb(th.accent)
            bar.line.fill.background()

        def footer(slide, n) -> None:
            strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Pt(6), W, Pt(6))
            strip.fill.solid(); strip.fill.fore_color.rgb = rgb(th.accent)
            strip.line.fill.background()
            textbox(slide, W - Inches(1.1), H - Inches(0.55), Inches(0.9), Inches(0.4),
                    str(n), size=11, color=th.dim, font=th.body_font, align=PP_ALIGN.RIGHT)

        # --- title slide: band + title + subtitle ---
        s0 = prs.slides.add_slide(blank)
        bg(s0)
        band = s0.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), W, Inches(2.9))
        band.fill.solid(); band.fill.fore_color.rgb = rgb(th.band)
        band.line.fill.background()
        textbox(s0, Inches(0.9), Inches(2.55), W - Inches(1.8), Inches(1.7),
                doc.title or deck.title or "Untitled",
                size=th.title_pt, color=th.band_ink, font=th.heading_font, bold=True)
        textbox(s0, Inches(0.9), Inches(4.15), W - Inches(1.8), Inches(0.6),
                f"{doc.kind.capitalize()}  ·  {len(deck.slides)} sections",
                size=14, color=th.band_ink, font=th.body_font)

        def head(s, title) -> None:
            textbox(s, Inches(0.9), Inches(0.55), W - Inches(1.8), Inches(1.0),
                    title, size=th.slide_title_pt, color=th.heading,
                    font=th.heading_font, bold=True)
            accent_line(s, Inches(1.5))

        def bullet_box(s, items, l, t, w, h, *, size=None) -> None:
            box = s.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame; tf.word_wrap = True
            for j, b in enumerate(items):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"•  {b}"; p.space_after = Pt(9)
                for r in p.runs:
                    r.font.name = th.body_font
                    r.font.size = Pt(size or th.body_pt)
                    r.font.color.rgb = rgb(th.ink)

        # section titles that have children (used for the divider layout)
        parent_titles = {top.title for top in doc.sections if top.children}
        imgs = list(self.images)

        # no attachments -> auto visuals on a subset of the plain bullet slides
        placeholder_slots: set[int] = set()
        if not imgs:
            standard = [
                idx for idx, sl in enumerate(deck.slides, 1)
                if self._layout(sl, sl.title in parent_titles, False) == "standard"
            ]
            placeholder_slots = set(standard[1::_PLACEHOLDER_EVERY])

        # --- content slides, with a layout chosen per slide ---
        for i, sl in enumerate(deck.slides, 1):
            b = [x for x in (sl.bullets or []) if x.strip()] or (
                [sl.notes.split(". ")[0]] if sl.notes else [])
            is_parent = sl.title in parent_titles
            lay = self._layout(sl, is_parent, bool(imgs) and not is_parent)
            img = imgs.pop(0) if (lay == "image" and imgs) else None
            auto_img = None
            if img is None and i in placeholder_slots:
                lay = "image"
                auto_img = _fetch_placeholder(
                    sl.title + " " + (b[0] if b else "")
                )

            s = prs.slides.add_slide(blank); bg(s)

            if lay == "divider":
                band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), W, Inches(2.3))
                band.fill.solid(); band.fill.fore_color.rgb = rgb(th.band)
                band.line.fill.background()
                textbox(s, Inches(0.9), Inches(3.0), W - Inches(1.8), Inches(1.6),
                        sl.title, size=th.title_pt - 4, color=th.band_ink,
                        font=th.heading_font, bold=True)

            elif lay == "quote":
                q = b[0].strip().strip('"“”')
                textbox(s, Inches(1.4), Inches(2.6), W - Inches(2.8), Inches(2.4),
                        f"“{q}”", size=26, color=th.heading,
                        font=th.heading_font, bold=False, align=PP_ALIGN.CENTER)
                textbox(s, Inches(1.4), Inches(5.1), W - Inches(2.8), Inches(0.5),
                        sl.title, size=13, color=th.dim, font=th.body_font,
                        align=PP_ALIGN.CENTER)

            elif lay == "bignum":
                head(s, sl.title)
                num = b[0].split()[0] if b else ""
                rest = b[0][len(num):].strip(" —-:·") if b else ""
                textbox(s, Inches(0.95), Inches(2.1), W - Inches(1.9), Inches(2.2),
                        num, size=84, color=th.accent, font=th.heading_font, bold=True)
                if rest:
                    textbox(s, Inches(0.98), Inches(4.4), W - Inches(1.9), Inches(0.8),
                            rest, size=20, color=th.ink, font=th.body_font)
                if len(b) > 1:
                    bullet_box(s, b[1:4], Inches(0.98), Inches(5.2), W - Inches(1.9),
                               Inches(1.6), size=15)

            elif lay == "twocol":
                head(s, sl.title)
                mid = (len(b) + 1) // 2
                bullet_box(s, b[:mid], Inches(0.95), Inches(1.95), Inches(5.9), Inches(4.6))
                bullet_box(s, b[mid:8], Inches(7.1), Inches(1.95), Inches(5.9), Inches(4.6))

            elif lay == "image":
                head(s, sl.title)
                pic = img or auto_img
                placed = False
                if pic:
                    try:
                        s.shapes.add_picture(pic, Inches(7.0), Inches(1.95),
                                             width=Inches(5.6))
                        placed = True
                    except Exception:  # noqa: BLE001 — a bad image shouldn't kill the deck
                        placed = False
                if not placed:
                    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0),
                                               Inches(1.95), Inches(5.6), Inches(3.95))
                    panel.fill.solid(); panel.fill.fore_color.rgb = rgb(th.band)
                    panel.line.color.rgb = rgb(th.accent)
                    textbox(s, Inches(7.3), Inches(3.5), Inches(5.0), Inches(1.0),
                            f"Visual — {sl.title}", size=13, color=th.band_ink,
                            font=th.body_font, align=PP_ALIGN.CENTER)
                    note = (s.notes_slide.notes_text_frame.text or "").strip()
                    s.notes_slide.notes_text_frame.text = (
                        (note + "\n\n" if note else "")
                        + f"[Add an image here: {sl.title}]"
                    )
                bullet_box(s, b[:6], Inches(0.95), Inches(1.95), Inches(5.7), Inches(4.6))

            else:  # standard
                head(s, sl.title)
                bullet_box(s, b[:7], Inches(0.95), Inches(1.95), W - Inches(1.9), Inches(4.7))

            footer(s, i)
            if sl.notes:
                s.notes_slide.notes_text_frame.text = sl.notes[:1500]

        # --- references slide ---
        cits = doc.all_citations()
        if cits:
            s = prs.slides.add_slide(blank)
            bg(s)
            textbox(s, Inches(0.9), Inches(0.55), W - Inches(1.8), Inches(1.0),
                    "References", size=th.slide_title_pt, color=th.heading,
                    font=th.heading_font, bold=True)
            accent_line(s, Inches(1.5))
            textbox(s, Inches(0.95), Inches(1.95), W - Inches(1.9), Inches(4.7),
                    [f"{i + 1}. {c.label or c.source or c.id}" for i, c in enumerate(cits)],
                    size=13, color=th.dim, font=th.body_font)

        buf = io.BytesIO()
        prs.save(buf)
        return RenderedDoc(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            text=_plain_preview(doc), ext="pptx", data=buf.getvalue(),
        )


class PdfRenderer:
    """Themed PDF via the bundled dependency-free writer (Helvetica, no install):
    an accent cover band, headings with an accent rule, wrapped body + bullets,
    a footer with the title and page number, a References page."""

    _M = 54.0        # page margin (pt)
    _PAGE_H = 792.0

    def __init__(self, theme=None) -> None:
        from app.services.authoring.theme import get_theme
        self.theme = theme if theme is not None else get_theme(None)

    def render(self, doc: DocumentModel) -> RenderedDoc:
        from app.services.authoring.pdf_writer import PDF

        th = self.theme
        p = PDF()
        W, M = p.w, self._M
        inner = W - 2 * M
        nums = doc.numbered_citations()
        self._page_i = 1

        def footer() -> None:
            p.line(M, self._PAGE_H - 40, W - M, self._PAGE_H - 40, th.dim, 0.5)
            p.text(M, self._PAGE_H - 34, (doc.title or "")[:70], size=8, color=th.dim)
            n = str(self._page_i)
            p.text(W - M - text_wf(n, 8), self._PAGE_H - 34, n, size=8, color=th.dim)

        def text_wf(s, size):
            from app.services.authoring.pdf_writer import text_width
            return text_width(s, size)

        def page_break(y):
            footer()
            p.new_page()
            self._page_i += 1
            return M + 24

        # --- cover ---
        p.rect(0, 0, W, 150, th.band)
        p.text(M, 54, (doc.title or "Untitled")[:64], size=28, bold=True, color=th.band_ink)
        p.text(M, 100, f"{doc.kind.capitalize()}  ·  {len(list(doc.walk()))} sections",
               size=11, color=th.band_ink)
        y = 200

        def emit(sec: Section, level: int, y: float) -> float:
            if y > self._PAGE_H - 110:
                y = page_break(y)
            hsize = {1: 17, 2: 14}.get(level, 12)
            p.text(M, y, sec.title, size=hsize, bold=True, color=th.heading)
            y += hsize + 4
            if level == 1:
                p.line(M, y, M + 120, y, th.accent, 2.0)
                y += 12
            for b in sec.blocks:
                if y > self._PAGE_H - 90:
                    y = page_break(y)
                cite = _num(b.citation_ids, nums)
                if b.kind in ("paragraph", "quote"):
                    y = p.paragraph(M + (14 if b.kind == "quote" else 0), y,
                                    (b.text + cite).strip(), size=11,
                                    color=th.dim if b.kind == "quote" else th.ink,
                                    max_w=inner - (14 if b.kind == "quote" else 0))
                    y += 6
                elif b.kind == "list":
                    for it in b.items:
                        if y > self._PAGE_H - 90:
                            y = page_break(y)
                        y = p.paragraph(M + 16, y, "•  " + it, size=11,
                                        color=th.ink, max_w=inner - 16)
                        y += 2
                    y += 4
                elif b.kind == "code":
                    for ln in b.text.split("\n")[:30]:
                        p.text(M + 8, y, ln[:100], size=9, color=th.ink)
                        y += 12
                    y += 4
                elif b.kind == "table" and b.rows:
                    for r_i, row in enumerate(b.rows[:20]):
                        p.text(M, y, "  |  ".join(str(c) for c in row)[:110],
                               size=10, bold=(r_i == 0), color=th.ink)
                        y += 14
                    y += 4
                elif b.kind == "figure":
                    y = p.paragraph(M, y, f"Figure: {b.caption or b.text}{cite}",
                                    size=10, color=th.dim, max_w=inner)
                    y += 6
            for child in sec.children:
                y = emit(child, level + 1, y + 6)
            return y + 10

        for sec in doc.sections:
            y = emit(sec, 1, y)

        cits = doc.all_citations()
        if cits:
            y = page_break(y)
            p.text(M, y, "References", size=17, bold=True, color=th.heading)
            y += 22
            p.line(M, y, M + 120, y, th.accent, 2.0)
            y += 14
            for i, c in enumerate(cits, 1):
                if y > self._PAGE_H - 90:
                    y = page_break(y)
                y = p.paragraph(M, y, f"{i}. {c.label or c.source or c.id}",
                                size=10, color=th.ink, max_w=inner)
                y += 4

        footer()
        return RenderedDoc("application/pdf", text=_plain_preview(doc), ext="pdf",
                           data=p.output())


_BY_NAME = {
    "markdown": MarkdownRenderer, "md": MarkdownRenderer,
    "html": HtmlRenderer,
    "docx": DocxRenderer, "pptx": PptxRenderer, "pdf": PdfRenderer,
}


def get_renderer(name: str, theme=None, images: list[str] | None = None) -> Renderer:
    cls = _BY_NAME.get(name.lower(), MarkdownRenderer)
    if cls in (MarkdownRenderer, HtmlRenderer):
        return cls()
    if cls is PptxRenderer:
        return PptxRenderer(theme=theme, images=images or [])
    return cls(theme=theme)
