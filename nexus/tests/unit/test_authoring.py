"""Acceptance (Unit): document model, renderers, outline, draft, review
(MILESTONE_M_PLAN.md §6)."""

from __future__ import annotations

import pytest

from app.llm.fake import ScriptedLLM
from app.services.authoring.draft import draft
from app.services.authoring.model import Block, Citation, DocumentModel, Section, SlideDeck
from app.services.authoring.outline import outline
from app.services.authoring.render import (
    DocxRenderer,
    HtmlRenderer,
    MarkdownRenderer,
    RendererUnavailable,
    get_renderer,
)
from app.services.authoring.review import review
from app.services.kb.store import KnowledgeBase


# --- model ---------------------------------------------------- #
def test_walk_and_citation_order() -> None:
    doc = DocumentModel(title="T")
    doc.add_citation(Citation(id="a", label="A"))
    doc.add_citation(Citation(id="b", label="B"))
    s1 = Section(title="One", blocks=[Block(kind="paragraph", text="x", citation_ids=["b"])])
    s2 = Section(title="Two", blocks=[Block(kind="paragraph", text="y", citation_ids=["a", "b"])])
    s1.children = [s2]
    doc.sections = [s1]
    assert [s.title for s in doc.walk()] == ["One", "Two"]
    assert [c.id for c in doc.all_citations()] == ["b", "a"]     # first-reference order
    assert doc.numbered_citations() == {"b": 1, "a": 2}


def test_slidedeck_one_per_h2() -> None:
    doc = DocumentModel(title="D")
    doc.sections = [
        Section(title="Intro", blocks=[Block(kind="list", items=["p1", "p2", "p3"])]),
        Section(title="Body", blocks=[Block(kind="paragraph", text="A full sentence. And more.")]),
    ]
    deck = SlideDeck.from_document(doc)
    assert [s.title for s in deck.slides] == ["Intro", "Body"]
    assert deck.slides[0].bullets == ["p1", "p2", "p3"]


# --- renderers -------------------------------------------- #
def _rich_doc() -> DocumentModel:
    doc = DocumentModel(title="Guide")
    doc.add_citation(Citation(id="c1", label="RFC 9110", source="https://rfc", trust="retrieved_web"))
    doc.sections = [Section(title="Overview", blocks=[
        Block(kind="paragraph", text="Intro text.", citation_ids=["c1"]),
        Block(kind="list", items=["one", "two"]),
        Block(kind="table", rows=[["h1", "h2"], ["a", "b"]]),
        Block(kind="code", lang="python", text="print(1)"),
        Block(kind="quote", text="a quote"),
    ])]
    return doc


def test_markdown_renderer_all_blocks_and_references() -> None:
    md = MarkdownRenderer().render(_rich_doc())
    assert md.mime == "text/markdown" and md.ext == "md"
    for frag in ("# Guide", "## Overview", "Intro text. [1]", "- one",
                 "| h1 | h2 |", "```python", "> a quote", "## References", "1. RFC 9110"):
        assert frag in md.text, frag


def test_html_renderer_is_nested_and_escaped() -> None:
    doc = _rich_doc()
    doc.sections[0].blocks[0].text = "a < b & c"
    h = HtmlRenderer().render(doc)
    assert "<article>" in h.text and "<section>" in h.text
    assert "a &lt; b &amp; c" in h.text
    assert '<section class="references">' in h.text


def test_docx_pptx_pdf_renderers_produce_valid_binaries() -> None:
    from app.services.authoring.render import PdfRenderer, PptxRenderer

    doc = _rich_doc()
    dx = DocxRenderer().render(doc)
    assert dx.ext == "docx" and dx.is_binary and dx.data[:2] == b"PK"  # zip/ooxml
    assert dx.text  # a readable preview is still attached
    px = PptxRenderer().render(doc)
    assert px.ext == "pptx" and px.is_binary and px.data[:2] == b"PK"
    pf = PdfRenderer().render(doc)
    assert pf.ext == "pdf" and pf.is_binary and pf.data[:8] == b"%PDF-1.4"
    assert pf.data.rstrip().endswith(b"%%EOF")

    # re-open to prove they are well-formed
    import io as _io

    import docx as _docx
    import pptx as _pptx
    from pypdf import PdfReader
    assert len(_docx.Document(_io.BytesIO(dx.data)).paragraphs) > 3
    assert len(_pptx.Presentation(_io.BytesIO(px.data)).slides._sldIdLst) >= 2
    assert len(PdfReader(_io.BytesIO(pf.data)).pages) >= 1

    assert isinstance(get_renderer("unknown"), MarkdownRenderer)


def test_pptx_auto_visuals_when_no_images_attached(monkeypatch) -> None:
    """With no attachments the deck still gets 'Visual —' placeholder panels on a
    subset of plain slides, speaker notes carry an [Add an image here] marker,
    and the file stays a valid OOXML package. NEXUS_DECK_IMAGES=0 forces the
    offline (panel) path so the test never hits the network."""
    import io as _io

    import pptx as _pptx

    from app.services.authoring.model import Block, DocumentModel, Section
    from app.services.authoring.render import PptxRenderer

    monkeypatch.setenv("NEXUS_DECK_IMAGES", "0")
    doc = DocumentModel(title="Deck")
    doc.sections = [
        Section(title=f"Point {n}", blocks=[Block(kind="list", items=["a", "b", "c"])])
        for n in range(1, 9)
    ]
    px = PptxRenderer().render(doc)  # no images= -> auto placeholders
    assert px.ext == "pptx" and px.data[:2] == b"PK"

    prs = _pptx.Presentation(_io.BytesIO(px.data))
    texts = [
        sh.text_frame.text
        for sl in prs.slides for sh in sl.shapes if sh.has_text_frame
    ]
    assert any(t.startswith("Visual —") for t in texts), "no placeholder panel rendered"
    notes = [
        sl.notes_slide.notes_text_frame.text
        for sl in prs.slides if sl.has_notes_slide
    ]
    assert any("[Add an image here" in n for n in notes)


# --- outline ------------------------------------------- #
def test_outline_builds_tree_and_flags_unsupported() -> None:
    kb = KnowledgeBase()
    kb.ingest_text("# Rate limit\n\nThe API allows 100 requests per minute.", uri="limits.md")
    llm = ScriptedLLM(['{"title": "API Guide", "sections": ['
                       '{"title": "Rate limits", "gist": "request quotas per minute"},'
                       '{"title": "Billing history", "gist": "invoices and refunds"}]}'])
    doc = outline("write an API guide", llm, kb=kb)
    assert doc.title == "API Guide"
    titles = [s.title for s in doc.walk()]
    assert "Rate limits" in titles
    billing = next(s for s in doc.walk() if s.title == "Billing history")
    assert "unsupported-section" in billing.flags
    kb.close()


def test_outline_empty_brief() -> None:
    doc = outline("", ScriptedLLM([]))
    assert doc.flags == ["empty-brief"] and len(doc.sections) == 1


# --- draft ----------------------------------------- #
def test_draft_fills_from_kb_with_citations() -> None:
    kb = KnowledgeBase()
    kb.ingest_text("# Rate limit\n\nThe API allows 100 requests per minute on free tier.", uri="l.md", title="L")

    def llm(system: str, prompt: str) -> str:
        s = system.lower()
        if "extract factual claims" in s:
            return '{"claims": [{"text": "100 requests per minute on free tier", "supported": true}]}'
        if "write one section" in s:
            return '{"paragraphs": [{"text": "The free tier permits 100 requests per minute.", "citation_ids": []}]}'
        return "{}"

    doc = DocumentModel(title="G")
    doc.sections = [Section(title="Rate limits", gist="request quotas"),
                    Section(title="Nonsense topic", gist="unrelated xyzzy")]
    draft(doc, ScriptedLLM(llm), brief="api guide", kb=kb)

    rl = doc.sections[0]
    assert any(b.citation_ids for b in rl.blocks)               # grounded
    assert doc.all_citations()[0].trust == "doc_input"
    ns = doc.sections[1]
    assert "unsupported-section" in ns.flags                    # no KB support
    kb.close()


# --- review --------------------------------------- #
def test_review_flags_structure_and_llm_issues() -> None:
    doc = DocumentModel(title="R")
    doc.sections = [
        Section(title="Empty", blocks=[]),
        Section(title="Empty", blocks=[Block(kind="paragraph", text="dup title section here.")]),
        Section(title="Factual", blocks=[Block(kind="paragraph",
                text="The system handles ten million requests. It never fails. Guaranteed.")]),
    ]
    llm = ScriptedLLM(['{"issues": [{"kind": "overclaim", "section": "Factual", '
                       '"detail": "never fails / guaranteed", "severity": "major"}]}'])
    issues = review(doc, llm)
    kinds = {i.kind for i in issues}
    assert {"empty-section", "duplicate-title", "missing-citation", "overclaim"} <= kinds
